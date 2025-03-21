#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Module de traitement des documents avec l'IA
Utilise Llama pour l'analyse et la personnalisation des documents
"""

import logging
import json
import os
import re
import chardet
import requests
import time
from typing import Dict, List, Optional, Tuple, Any

logging.basicConfig(level=logging.DEBUG)

logger = logging.getLogger("VynalDocsAutomator.AIDocumentProcessor")

class AIDocumentProcessor:
    """
    Classe pour le traitement des documents avec l'IA
    Utilise Llama pour analyser et personnaliser les documents
    """
    
    def __init__(self):
        """
        Initialise le processeur de documents
        """
        self.config = self._load_config()
        self.cache = {}
        self.cache_size = 100  # Réduit de 1000 à 100
        self.fallback_mode = False  # Ajout de l'attribut fallback_mode initialisé à False
        
    def _load_config(self):
        """
        Charge la configuration depuis le fichier config.json
        Gère les différents encodages pour éviter les problèmes de décodage
        """
        try:
            # Utiliser le chemin absolu pour éviter les problèmes de chemin relatif
            config_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "config", "ai_config.json")
            logger.info(f"Tentative de chargement de la configuration depuis: {config_path}")
            
            if not os.path.exists(config_path):
                logger.warning(f"Le fichier de configuration n'existe pas: {config_path}")
                # Si le fichier n'existe pas, retourner une configuration par défaut
                return {
                    "api_url": "http://localhost:11434/api/generate",
                    "model": "mistral:latest",
                    "options": {
                        "temperature": 0.7,
                        "max_tokens": 512
                    }
                }
            
            # Essayer différents encodages
            encodings = ['utf-8', 'latin-1', 'cp1252', 'utf-16', 'utf-16-le', 'utf-16-be']
            
            for encoding in encodings:
                try:
                    with open(config_path, 'r', encoding=encoding) as f:
                        config = json.load(f)
                    logger.info(f"Configuration chargée avec l'encodage {encoding} depuis {config_path}")
                    return config
                except UnicodeDecodeError:
                    logger.debug(f"Échec du décodage avec l'encodage {encoding}")
                    continue  # Essayer le prochain encodage
                except json.JSONDecodeError:
                    # Si le fichier s'ouvre mais que le JSON est invalide
                    logger.warning(f"Format JSON invalide dans le fichier de configuration avec l'encodage {encoding}")
                    continue
            
            # Si on arrive ici, tous les encodages ont échoué
            # Essai en mode binaire avec détection d'encodage
            logger.warning("Tous les encodages standards ont échoué, tentative de détection automatique")
            import chardet
            try:
                with open(config_path, 'rb') as f:
                    raw_data = f.read()
                    result = chardet.detect(raw_data)
                    detected_encoding = result['encoding']
                    logger.info(f"Encodage détecté: {detected_encoding} (confiance: {result['confidence']})")
                    
                    if detected_encoding:
                        try:
                            config_str = raw_data.decode(detected_encoding)
                            config = json.loads(config_str)
                            logger.info(f"Configuration chargée avec l'encodage détecté {detected_encoding}")
                            return config
                        except Exception as e:
                            logger.warning(f"Échec avec l'encodage détecté {detected_encoding}: {e}")
            except Exception as e:
                logger.warning(f"Échec de la détection d'encodage: {e}")
                
            # Créer une nouvelle configuration par défaut
            logger.warning("Création d'une nouvelle configuration par défaut")
            config = {
                "api_url": "http://localhost:11434/api/generate",
                "model": "mistral:latest",
                "options": {
                    "temperature": 0.7,
                    "max_tokens": 512
                }
            }
            
            # Essayer de sauvegarder cette configuration
            try:
                os.makedirs(os.path.dirname(config_path), exist_ok=True)
                with open(config_path, 'w', encoding='utf-8') as f:
                    json.dump(config, f, indent=4)
                logger.info("Nouvelle configuration par défaut créée")
            except Exception as e:
                logger.error(f"Impossible de créer la configuration par défaut: {e}")
                
            return config
            
        except Exception as e:
            logger.warning(f"Impossible de charger la configuration: {e}")
            # Configuration par défaut
            return {
                "api_url": "http://localhost:11434/api/generate",
                "model": "mistral:latest",
                "options": {
                    "temperature": 0.7,
                    "max_tokens": 512
                }
            }
        
    def _read_file_safely(self, file_path: str) -> str:
        """
        Lit un fichier en détectant automatiquement l'encodage
        
        Args:
            file_path: Chemin vers le fichier à lire
            
        Returns:
            Contenu du fichier
        """
        try:
            # Détecter le type de fichier par l'extension
            _, ext = os.path.splitext(file_path)
            ext = ext.lower()
            
            # Fichiers PDF
            if ext == '.pdf':
                try:
                    import PyPDF2
                    with open(file_path, 'rb') as f:
                        reader = PyPDF2.PdfReader(f)
                        return " ".join([page.extract_text() for page in reader.pages])
                except ImportError:
                    logger.warning("PyPDF2 non disponible, traitement du PDF en tant que fichier binaire")
                    return f"[PDF NON SUPPORTÉ: {os.path.basename(file_path)}]"
                except Exception as e:
                    logger.error(f"Erreur lors de la lecture du PDF: {e}")
                    return f"[ERREUR PDF: {os.path.basename(file_path)}]"
            
            # Fichiers DOCX
            elif ext == '.docx':
                try:
                    import docx
                    doc = docx.Document(file_path)
                    return " ".join([para.text for para in doc.paragraphs])
                except ImportError:
                    logger.warning("python-docx non disponible, traitement du DOCX en tant que fichier binaire")
                    return f"[DOCX NON SUPPORTÉ: {os.path.basename(file_path)}]"
                except Exception as e:
                    logger.error(f"Erreur lors de la lecture du DOCX: {e}")
                    return f"[ERREUR DOCX: {os.path.basename(file_path)}]"
            
            # Fichiers Excel
            elif ext in ['.xlsx', '.xls']:
                try:
                    import pandas as pd
                    df = pd.read_excel(file_path)
                    return df.to_string()
                except ImportError:
                    logger.warning("pandas non disponible, traitement du fichier Excel en tant que fichier binaire")
                    return f"[EXCEL NON SUPPORTÉ: {os.path.basename(file_path)}]"
                except Exception as e:
                    logger.error(f"Erreur lors de la lecture du fichier Excel: {e}")
                    return f"[ERREUR EXCEL: {os.path.basename(file_path)}]"
            
            # Fichiers PowerPoint
            elif ext in ['.pptx', '.ppt']:
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text.append(shape.text)
                    return "\n".join(text)
                except ImportError:
                    logger.warning("python-pptx non disponible, traitement du fichier PowerPoint en tant que fichier binaire")
                    return f"[POWERPOINT NON SUPPORTÉ: {os.path.basename(file_path)}]"
                except Exception as e:
                    logger.error(f"Erreur lors de la lecture du fichier PowerPoint: {e}")
                    return f"[ERREUR POWERPOINT: {os.path.basename(file_path)}]"
            
            # Fichiers OpenDocument
            elif ext in ['.odt', '.ods', '.odp']:
                try:
                    import odf
                    if ext == '.odt':
                        doc = odf.text.TextDocument(file_path)
                        return " ".join([p.text for p in doc.getElementsByType(odf.text.P)])
                    elif ext == '.ods':
                        doc = odf.spreadsheet.SpreadsheetDocument(file_path)
                        return " ".join([cell.text for table in doc.getElementsByType(odf.table.Table) for cell in table.getElementsByType(odf.table.Cell)])
                    else:  # .odp
                        doc = odf.presentation.PresentationDocument(file_path)
                        return " ".join([shape.text for page in doc.getElementsByType(odf.presentation.Page) for shape in page.getElementsByType(odf.presentation.Text)])
                except ImportError:
                    logger.warning("odfpy non disponible, traitement du fichier OpenDocument en tant que fichier binaire")
                    return f"[OPENDOCUMENT NON SUPPORTÉ: {os.path.basename(file_path)}]"
                except Exception as e:
                    logger.error(f"Erreur lors de la lecture du fichier OpenDocument: {e}")
                    return f"[ERREUR OPENDOCUMENT: {os.path.basename(file_path)}]"
            
            # Fichiers RTF
            elif ext == '.rtf':
                try:
                    import striprtf
                    with open(file_path, 'r', encoding='utf-8') as f:
                        rtf_content = f.read()
                    return striprtf.rtf_to_text(rtf_content)
                except ImportError:
                    logger.warning("striprtf non disponible, traitement du fichier RTF en tant que fichier binaire")
                    return f"[RTF NON SUPPORTÉ: {os.path.basename(file_path)}]"
                except Exception as e:
                    logger.error(f"Erreur lors de la lecture du fichier RTF: {e}")
                    return f"[ERREUR RTF: {os.path.basename(file_path)}]"
            
            # Fichiers HTML
            elif ext in ['.html', '.htm']:
                try:
                    from bs4 import BeautifulSoup
                    with open(file_path, 'r', encoding='utf-8') as f:
                        soup = BeautifulSoup(f.read(), 'html.parser')
                    return soup.get_text()
                except ImportError:
                    logger.warning("beautifulsoup4 non disponible, traitement du fichier HTML en tant que fichier binaire")
                    return f"[HTML NON SUPPORTÉ: {os.path.basename(file_path)}]"
                except Exception as e:
                    logger.error(f"Erreur lors de la lecture du fichier HTML: {e}")
                    return f"[ERREUR HTML: {os.path.basename(file_path)}]"
            
            # Fichiers XML
            elif ext == '.xml':
                try:
                    import xml.etree.ElementTree as ET
                    tree = ET.parse(file_path)
                    root = tree.getroot()
                    return ET.tostring(root, encoding='unicode', method='text')
                except Exception as e:
                    logger.error(f"Erreur lors de la lecture du fichier XML: {e}")
                    return f"[ERREUR XML: {os.path.basename(file_path)}]"
            
            # Fichiers JSON
            elif ext == '.json':
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        data = json.load(f)
                    return json.dumps(data, ensure_ascii=False, indent=2)
                except Exception as e:
                    logger.error(f"Erreur lors de la lecture du fichier JSON: {e}")
                    return f"[ERREUR JSON: {os.path.basename(file_path)}]"
            
            # Fichiers CSV
            elif ext == '.csv':
                try:
                    import pandas as pd
                    df = pd.read_csv(file_path)
                    return df.to_string()
                except ImportError:
                    logger.warning("pandas non disponible, traitement du fichier CSV en tant que fichier binaire")
                    return f"[CSV NON SUPPORTÉ: {os.path.basename(file_path)}]"
                except Exception as e:
                    logger.error(f"Erreur lors de la lecture du fichier CSV: {e}")
                    return f"[ERREUR CSV: {os.path.basename(file_path)}]"
            
            # Fichiers Markdown
            elif ext in ['.md', '.markdown']:
                try:
                    import markdown
                    with open(file_path, 'r', encoding='utf-8') as f:
                        md_content = f.read()
                    return markdown.markdown(md_content)
                except ImportError:
                    logger.warning("markdown non disponible, traitement du fichier Markdown en tant que fichier binaire")
                    return f"[MARKDOWN NON SUPPORTÉ: {os.path.basename(file_path)}]"
                except Exception as e:
                    logger.error(f"Erreur lors de la lecture du fichier Markdown: {e}")
                    return f"[ERREUR MARKDOWN: {os.path.basename(file_path)}]"
            
            # Images
            elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
                try:
                    import pytesseract
                    from PIL import Image
                    img = Image.open(file_path)
                    return pytesseract.image_to_string(img)
                except ImportError:
                    logger.warning("pytesseract non disponible, traitement de l'image en tant que fichier binaire")
                    return f"[IMAGE NON SUPPORTÉ: {os.path.basename(file_path)}]"
                except Exception as e:
                    logger.error(f"Erreur lors de la lecture de l'image: {e}")
                    return f"[ERREUR IMAGE: {os.path.basename(file_path)}]"
            
            # Fichiers texte - détecter l'encodage automatiquement
            else:
                with open(file_path, 'rb') as f:
                    raw_data = f.read()
                    result = chardet.detect(raw_data)
                    encoding = result['encoding'] if result['encoding'] else 'utf-8'
                    
                try:
                    return raw_data.decode(encoding)
                except UnicodeDecodeError:
                    # En cas d'échec, essayer latin-1 qui ne peut jamais échouer
                    return raw_data.decode('latin-1')
                    
        except Exception as e:
            logger.error(f"Erreur lors de la lecture du fichier: {e}")
            return f"[ERREUR DE LECTURE: {os.path.basename(file_path)}]"
            
    def analyze_document(self, file_path: str, file_type: str = None, encoding: str = 'utf-8') -> Dict[str, Any]:
        """
        Analyse un document pour extraire des informations pertinentes
        
        Args:
            file_path: Chemin vers le document à analyser
            file_type: Type du fichier (extension sans le point)
            encoding: Encodage du fichier (par défaut utf-8)
            
        Returns:
            Dict contenant les informations extraites
        """
        if not os.path.exists(file_path):
            return {"error": f"Le fichier {file_path} n'existe pas"}
            
        if not file_type:
            # Extraire l'extension du fichier
            _, ext = os.path.splitext(file_path)
            file_type = ext.lower()[1:] if ext else ""
            
        # Convertir le document en texte
        try:
            content = self.convert_to_text(file_path, file_type, encoding)
        except Exception as e:
            logger.error(f"Erreur lors de la conversion du document en texte: {e}")
            return {"error": f"Impossible de lire le document: {str(e)}"}
            
        if not content or content.startswith("Erreur:"):
            return {"error": content if content else "Contenu vide ou non lisible"}
            
        # Utiliser directement le mode de secours si requis
        if self.fallback_mode:
            logger.info("Mode de secours activé, utilisation directe du mode de secours")
            return self._analyze_fallback(content)
            
        # Analyser le document
        try:
            # Analyse par sections pour les documents longs
            result = self._analyze_document_by_sections(content)
            
            # Si résultat vide ou erreur, utiliser le mode de secours
            if not result or "error" in result or not result.get("variables"):
                logger.warning("Analyse par sections a échoué ou n'a trouvé aucune variable, tentative avec le mode de secours")
                result = self._analyze_fallback(content)
            
            logger.info(f"Analyse terminée avec succès: {len(result['variables'])} variables trouvées")
            return result
            
        except Exception as e:
            logger.error(f"Erreur lors de l'analyse du document: {e}")
            return self._analyze_fallback(content)
    
    def _smart_reduce_content(self, content: str, max_size: int) -> str:
        """
        Réduit intelligemment le contenu pour l'analyse en préservant les parties importantes
        
        Args:
            content: Contenu à réduire
            max_size: Taille maximale souhaitée
            
        Returns:
            Contenu réduit intelligemment
        """
        if len(content) <= max_size:
            return content
            
        # Zones prioritaires
        start_size = min(max_size // 3, 800)  # Début du document (en-têtes, intro)
        end_size = min(max_size // 3, 600)    # Fin du document (conclusion, signatures)
        
        # Extraire les sections importantes
        content_start = content[:start_size]
        content_end = content[-end_size:]
        
        # Espace restant pour les sections importantes
        remaining_size = max_size - (start_size + end_size + 100)  # 100 pour le texte de séparation
        
        # Mots-clés importants selon le type de document
        important_keywords = [
            # Identités
            "client", "nom", "prénom", "société", "entreprise", 
            # Coordonnées
            "adresse", "email", "téléphone", "portable", "contact",
            # Données financières 
            "montant", "prix", "total", "somme", "€", "euros", "eur",
            # Références et dates
            "référence", "facture", "commande", "date", "numéro",
            # Termes légaux
            "contrat", "signature", "légal", "obligation"
        ]
        
        # Découper en paragraphes et noter leur importance
        paragraphs = re.split(r'\n\s*\n', content)
        scored_paragraphs = []
        
        for para in paragraphs:
            # Ignorer les paragraphes vides ou trop courts
            if len(para.strip()) < 5:
                continue
                
            # Calculer un score d'importance
            score = 0
            para_lower = para.lower()
            
            # Points pour chaque mot-clé trouvé
            for keyword in important_keywords:
                if keyword in para_lower:
                    score += 2
                    
            # Points supplémentaires pour les formats spéciaux
            # Dates
            if re.search(r'\b\d{1,2}[/.-]\d{1,2}[/.-]\d{2,4}\b', para):
                score += 3
            # Montants
            if re.search(r'\b\d+(?:[,.]\d{1,2})?(?:\s?[€$]|EUR)?\b', para):
                score += 3
            # Emails
            if re.search(r'\b[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+\b', para):
                score += 3
            # Téléphones
            if re.search(r'\b(?:0\d[\s.-]?){4,5}\d{2}\b', para):
                score += 3
            
            # Ajouter le paragraphe avec son score
            scored_paragraphs.append((para, score))
        
        # Trier les paragraphes par score d'importance décroissant
        scored_paragraphs.sort(key=lambda x: x[1], reverse=True)
        
        # Sélectionner les paragraphes les plus importants
        selected_paragraphs = []
        total_length = 0
        
        for para, score in scored_paragraphs:
            if total_length + len(para) <= remaining_size:
                selected_paragraphs.append(para)
                total_length += len(para) + 2  # +2 pour les sauts de ligne
            else:
                break
        
        # Construire le contenu final
        reduced_content = content_start + "\n\n"
        if selected_paragraphs:
            reduced_content += "--- INFORMATIONS IMPORTANTES ---\n\n" + "\n\n".join(selected_paragraphs) + "\n\n"
        reduced_content += "--- FIN DU DOCUMENT ---\n\n" + content_end
        
        logger.info(f"Contenu réduit de {len(content)} à {len(reduced_content)} caractères")
        return reduced_content

    def convert_to_text(self, file_path, file_type=None, encoding='utf-8'):
        """
        Convertit un document en texte brut, quel que soit son format
        
        Args:
            file_path: Chemin du fichier à convertir
            file_type: Type du fichier (extension)
            encoding: Encodage du fichier
            
        Returns:
            str: Contenu textuel du document
        """
        if not file_type:
            file_type = os.path.splitext(file_path)[1].lower()[1:]
            
        try:
            # Si le fichier est un .docx (Word)
            if file_type in ['docx', 'doc']:
                try:
                    import docx
                    doc = docx.Document(file_path)
                    full_text = []
                    # Extraire le texte de chaque paragraphe
                    for para in doc.paragraphs:
                        full_text.append(para.text)
                    # Extraire le texte des tableaux
                    for table in doc.tables:
                        for row in table.rows:
                            for cell in row.cells:
                                for para in cell.paragraphs:
                                    full_text.append(para.text)
                    return '\n'.join(full_text)
                except ImportError:
                    logger.warning("Module python-docx non disponible, tentative de fallback")
                    # Fallback à une lecture textuelle
                    try:
                        with open(file_path, 'rb') as f:
                            import chardet
                            result = chardet.detect(f.read())
                            detected_encoding = result['encoding']
                        with open(file_path, 'r', encoding=detected_encoding or 'utf-8', errors='replace') as f:
                            return f.read()
                    except Exception as e:
                        logger.error(f"Erreur lors de la lecture du fichier Word: {e}")
                        return f"Erreur: {str(e)}"
                except Exception as e:
                    logger.error(f"Erreur lors de l'extraction du texte du fichier Word: {e}")
                    return f"Erreur: {str(e)}"
                    
            # Si le fichier est un .pdf
            elif file_type == 'pdf':
                try:
                    import PyPDF2
                    with open(file_path, 'rb') as f:
                        pdf_reader = PyPDF2.PdfReader(f)
                        text = ''
                        for page_num in range(len(pdf_reader.pages)):
                            page = pdf_reader.pages[page_num]
                            text += page.extract_text() + '\n'
                        return text
                except ImportError:
                    logger.warning("Module PyPDF2 non disponible, tentative de fallback")
                    # Fallback vers une lecture textuelle basique
                    try:
                        with open(file_path, 'r', encoding=encoding, errors='replace') as f:
                            return f.read()
                    except Exception as e:
                        logger.error(f"Erreur lors de la lecture du fichier PDF: {e}")
                        return f"Erreur: {str(e)}"
                except Exception as e:
                    logger.error(f"Erreur lors de l'extraction du texte du fichier PDF: {e}")
                    return f"Erreur: {str(e)}"
                    
            # Si le fichier est une image
            elif file_type in ['jpg', 'jpeg', 'png', 'gif', 'bmp']:
                try:
                    from PIL import Image
                    import pytesseract
                    image = Image.open(file_path)
                    text = pytesseract.image_to_string(image, lang='fra+eng')
                    return text
                except ImportError:
                    logger.warning("Module pytesseract non disponible")
                    return "Erreur: L'extraction de texte depuis des images nécessite le module pytesseract"
                except Exception as e:
                    logger.error(f"Erreur lors de l'OCR sur l'image: {e}")
                    return f"Erreur: {str(e)}"
                    
            # Pour les fichiers texte (.txt, .md, .html, etc.)
            else:
                # Essayer différents encodages pour gérer tous les cas de figure
                encodings = ['utf-8', 'latin-1', 'cp1252', 'utf-16', 'utf-16-le', 'utf-16-be']
                content = None
                last_error = None
                
                # Si un encodage spécifique est demandé, le tester en premier
                if encoding and encoding not in encodings:
                    encodings.insert(0, encoding)
                
                for enc in encodings:
                    try:
                        with open(file_path, 'r', encoding=enc) as f:
                            content = f.read()
                        logger.info(f"Fichier lu avec succès en utilisant l'encodage {enc}")
                        break
                    except UnicodeDecodeError as e:
                        logger.debug(f"Échec de lecture avec l'encodage {enc}: {e}")
                        last_error = e
                    except Exception as e:
                        logger.warning(f"Erreur lors de la lecture du fichier avec l'encodage {enc}: {e}")
                        last_error = e
                
                # Si tous les encodages ont échoué, tenter une détection automatique
                if content is None:
                    try:
                        import chardet
                        with open(file_path, 'rb') as f:
                            raw_data = f.read()
                            result = chardet.detect(raw_data)
                            detected_encoding = result['encoding']
                            
                            if detected_encoding:
                                logger.info(f"Encodage détecté: {detected_encoding}")
                                with open(file_path, 'r', encoding=detected_encoding, errors='replace') as f:
                                    content = f.read()
                            else:
                                # Dernier recours: lecture en mode binaire et conversion forcée
                                content = raw_data.decode('utf-8', errors='replace')
                    except Exception as e:
                        logger.error(f"Erreur lors de la détection de l'encodage: {e}")
                        if last_error:
                            raise last_error
                        else:
                            raise e
                
                return content
                
        except Exception as e:
            logger.error(f"Erreur lors de la conversion du document en texte: {e}")
            return f"Erreur lors de la lecture du document: {str(e)}"
            
    def fill_template(self, template_path: str, client_info: Dict, encoding: str = 'utf-8') -> Tuple[str, Dict]:
        """
        Remplit un modèle avec les informations client
        
        Args:
            template_path: Chemin vers le modèle
            client_info: Informations client
            encoding: Encodage du fichier
            
        Returns:
            Tuple[str, Dict]: (contenu rempli, variables manquantes)
        """
        try:
            logger.info(f"Début du remplissage du modèle: {template_path}")
            
            # 1. Validation initiale
            if not os.path.exists(template_path):
                raise FileNotFoundError(f"Le modèle n'existe pas: {template_path}")
            
            if not client_info or not isinstance(client_info, dict):
                raise ValueError("Les informations client sont invalides")
            
            # 2. Lecture du modèle avec gestion d'encodage
            content = None
            encoding_errors = []
            
            for enc in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    with open(template_path, 'r', encoding=enc) as f:
                        content = f.read()
                        logger.debug(f"Lecture réussie avec l'encodage: {enc}")
                        break
                except UnicodeDecodeError as e:
                    encoding_errors.append(f"Erreur avec {enc}: {str(e)}")
                    continue
                    
            if content is None:
                raise UnicodeError(f"Impossible de lire le fichier. Erreurs: {', '.join(encoding_errors)}")
            
            # 3. Analyse du modèle
            logger.debug("Analyse du modèle...")
            variables = self._analyze_template_content(content)
            
            # 4. Validation des variables requises
            missing_vars = {}
            invalid_vars = []
            
            for var_name, var_info in variables.items():
                if var_info.get('required', False):
                    if var_name not in client_info:
                        missing_vars[var_name] = var_info
                    else:
                        # Valider la valeur selon le type
                        value = client_info[var_name]
                        validation_rules = self._get_validation_rules(var_info.get('type', 'text'))
                        if not self._validate_variable(value, validation_rules):
                            invalid_vars.append(var_name)
            
            if invalid_vars:
                logger.warning(f"Variables invalides: {invalid_vars}")
                raise ValueError(f"Les variables suivantes sont invalides: {', '.join(invalid_vars)}")
            
            # 5. Remplissage du modèle
            if not self.fallback_mode and len(missing_vars) == 0:
                try:
                    logger.info("Utilisation de l'IA pour le remplissage...")
                    filled_content = self._fill_with_ai(content, client_info)
                except Exception as e:
                    logger.error(f"Erreur lors du remplissage IA: {e}")
                    filled_content = self._fill_template_basic(content, client_info)
                else:
                    logger.info("Utilisation du remplissage basique...")
                    filled_content = self._fill_template_basic(content, client_info)
            
            return filled_content, missing_vars
            
        except Exception as e:
            logger.error(f"Erreur lors du remplissage du modèle: {e}")
            return str(e), {
                    "error": {
                        "type": "text",
                    "description": "Erreur de traitement",
                        "required": True,
                    "current_value": str(e)
                }
            }
    
    def _validate_variable(self, value: Any, rules: Dict[str, Any]) -> bool:
        """
        Valide une variable selon les règles définies
        
        Args:
            value: Valeur à valider
            rules: Règles de validation
            
        Returns:
            bool: True si la valeur est valide
        """
        try:
            if 'regex' in rules:
                if not re.match(rules['regex'], str(value)):
                    return False
            
            if 'min_length' in rules and isinstance(value, (str, list)):
                if len(value) < rules['min_length']:
                    return False
            
            if 'max_length' in rules and isinstance(value, (str, list)):
                if len(value) > rules['max_length']:
                    return False
            
            if 'min_value' in rules and isinstance(value, (int, float)):
                if value < rules['min_value']:
                    return False
            
            if 'decimal_places' in rules and isinstance(value, (int, float)):
                str_value = str(value)
                if '.' in str_value:
                    if len(str_value.split('.')[1]) > rules['decimal_places']:
                        return False
            
            return True
            
        except Exception as e:
            logger.error(f"Erreur lors de la validation: {e}")
            return False

    def generate_final_document(self, template_path: str, all_variables: Dict[str, str]) -> str:
        """
        Génère le document final avec toutes les variables
        
        Args:
            template_path: Chemin vers le modèle
            all_variables: Toutes les variables (client + complétées)
            
        Returns:
            Texte final du document
        """
        try:
            # Lire le contenu du modèle
            content = self._read_file_safely(template_path)
            original_content = content
            
            # Tentative initiale de remplacement simple des variables
            for var_name, value in all_variables.items():
                # S'assurer que la valeur est une chaîne non-vide
                value_str = str(value) if value is not None else ""
                
                # Différents formats de variables à remplacer
                patterns = [
                    f"{{{var_name}}}",     # {variable}
                    f"[{var_name}]",       # [variable]
                    f"<{var_name}>",       # <variable>
                    f"${var_name}$",       # $variable$
                    f"«{var_name}»"        # «variable»
                ]
                
                for pattern in patterns:
                    content = content.replace(pattern, value_str)
            
            # Vérifier si des variables n'ont pas été remplacées (motifs de variables restants)
            remaining_vars = []
            for pattern in [r'\{[a-zA-Z0-9_]+\}', r'\[[a-zA-Z0-9_]+\]', r'<[a-zA-Z0-9_]+>', r'\$[a-zA-Z0-9_]+\$', r'«[a-zA-Z0-9_]+»']:
                matches = re.findall(pattern, content)
                remaining_vars.extend(matches)
            
            # Si des variables n'ont pas été remplacées, utiliser l'IA pour un remplacement plus intelligent
            if remaining_vars and not self.fallback_mode:
                logger.info(f"{len(remaining_vars)} variables non remplacées détectées, utilisation de l'IA pour amélioration")
                
                # Utiliser l'IA pour générer un document complet
                final_content = self._complete_document_with_ai(content, all_variables, remaining_vars)
                
                # Vérification de qualité du document final
                if self._verify_document_quality(final_content, original_content, all_variables):
                    logger.info("Document final généré avec succès via IA")
                    return final_content
                else:
                    logger.warning("Document généré par IA de qualité insuffisante, fallback vers remplacement simple")
            
            # Si aucune variable non remplacée ou si en mode fallback, retourner le document modifié
            logger.info("Document final généré avec succès via remplacement simple")
            return content
            
        except Exception as e:
            logger.error(f"Erreur lors de la génération du document final: {e}")
            
            # Tentative de récupération en cas d'erreur
            try:
                # Lecture du template original et remplacement basique
                with open(template_path, 'r', encoding='utf-8', errors='replace') as f:
                    recovery_content = f.read()
                
                # Remplacement forcé avec toutes les variables
                for var_name, value in all_variables.items():
                    value_str = str(value) if value is not None else ""
                    for pattern in [f"{{{var_name}}}", f"[{var_name}]", f"<{var_name}>"]:
                        recovery_content = recovery_content.replace(pattern, value_str)
                
                logger.info("Document récupéré avec succès en mode de secours")
                return recovery_content
            except Exception as recovery_error:
                logger.error(f"Échec de la récupération: {recovery_error}")
                
                # En dernier recours, générer un document minimal mais complet
                minimal_doc = f"""
                # Document généré pour {all_variables.get('name', 'Client')}
                
                Date: {all_variables.get('date', '')}
                Référence: {all_variables.get('reference', '')}
                
                """
                
                # Ajouter toutes les variables disponibles
                for var_name, value in all_variables.items():
                    if var_name not in ['name', 'date', 'reference']:
                        minimal_doc += f"{var_name}: {value}\n"
                
                logger.warning("Document minimal généré suite à une erreur critique")
                return minimal_doc
    
    def _complete_document_with_ai(self, content: str, variables: Dict[str, str], remaining_vars: List[str]) -> str:
        """
        Utilise l'IA pour compléter intelligemment le document avec toutes les variables
        
        Args:
            content: Contenu du document partiellement rempli
            variables: Toutes les variables disponibles
            remaining_vars: Liste des variables non remplacées
            
        Returns:
            Document complété par l'IA
        """
        # Formater pour plus de lisibilité
        remaining_formatted = ', '.join(remaining_vars[:10])
        if len(remaining_vars) > 10:
            remaining_formatted += f" et {len(remaining_vars) - 10} autres"
        
        # Préparation du prompt
        prompt = f"""Vous êtes un expert en génération de documents juridiques et administratifs.

DOCUMENT ACTUEL:
```
{content}
```

VARIABLES DISPONIBLES:
```
{json.dumps(variables, indent=2, ensure_ascii=False)}
```

VARIABLES NON REMPLACÉES: {remaining_formatted}

TÂCHE:
1. Analysez le document et identifiez toutes les variables non remplacées
2. Remplacez chaque variable par sa valeur correspondante dans le dictionnaire
3. Si une variable n'a pas de valeur correspondante, utilisez une valeur standard appropriée
4. Assurez-vous que le texte final est cohérent et complet
5. Conservez EXACTEMENT la mise en forme du document (espaces, retours à la ligne, etc.)

IMPORTANT:
- Ne retournez que le document complété, sans explication
- Respectez la structure exacte du document original
- Utilisez toujours les informations fournies dans le dictionnaire quand elles existent

DOCUMENT COMPLÉTÉ:
"""
        
        try:
            # Appel à Ollama avec un timeout plus long pour garantir l'achèvement
            response = self._call_ollama(prompt, timeout=180)
            
            # Vérifier si la réponse contient du texte
            if not response or len(response) < 50:
                logger.warning("Réponse de l'IA trop courte ou vide")
                return content
                
            return response
            
        except Exception as e:
            logger.error(f"Échec de la complétion IA: {e}")
            return content
    
    def _verify_document_quality(self, final_doc: str, original_doc: str, variables: Dict[str, str]) -> bool:
        """
        Vérifie la qualité du document généré pour s'assurer qu'il est complet
        
        Args:
            final_doc: Document généré
            original_doc: Document original
            variables: Variables utilisées
            
        Returns:
            True si le document est de bonne qualité, False sinon
        """
        # 1. Vérifier que le document n'est pas vide
        if not final_doc or len(final_doc) < 50:
            logger.warning("Document final trop court")
            return False
            
        # 2. Vérifier que le document a une taille raisonnable par rapport à l'original
        if len(final_doc) < len(original_doc) * 0.7:
            logger.warning(f"Document final trop court: {len(final_doc)} vs {len(original_doc)} caractères")
            return False
            
        # 3. Vérifier que les variables importantes sont présentes dans le document final
        important_vars = ['name', 'date', 'montant', 'reference', 'adresse', 'telephone']
        for var in important_vars:
            if var in variables and variables[var]:
                if variables[var] not in final_doc:
                    logger.warning(f"Variable importante '{var}' absente du document final")
                    return False
        
        # 4. Vérifier qu'il ne reste pas de motifs de variables
        for pattern in [r'\{[a-zA-Z0-9_]+\}', r'\[[a-zA-Z0-9_]+\]', r'<[a-zA-Z0-9_]+>', r'\$[a-zA-Z0-9_]+\$', r'«[a-zA-Z0-9_]+»']:
            if re.search(pattern, final_doc):
                logger.warning(f"Variables non remplacées trouvées dans le document final")
                return False
                
        return True

    def personalize_document(self, content: str, variables: Dict[str, Any]) -> str:
        """
        Personnalise un document en remplaçant les variables par leurs valeurs
        
        Args:
            content: Contenu du document à personnaliser
            variables: Dictionnaire des variables à remplacer
            
        Returns:
            Document personnalisé
        """
        try:
            # Préparer les variables
            vars_list = []
            for var_name, var_info in variables.items():
                value = ""
                if isinstance(var_info, dict) and "current_value" in var_info:
                    value = var_info["current_value"]
                    description = var_info.get("description", var_name)
                    var_type = var_info.get("type", "text")
                    vars_list.append(f"{var_name} ({description}, {var_type}): {value}")
                else:
                    value = str(var_info)
                    vars_list.append(f"{var_name}: {value}")
                    
            vars_json = "\n".join(vars_list)
            
            # Créer le prompt pour personnaliser le document
            prompt = f"""
Personnalise ce document en remplaçant les variables par les valeurs données:

Document:
```
{content}
```

Variables:
{vars_json}

Retourne uniquement le document personnalisé."""

            # Si le document est petit, le traiter en une fois
            if len(content) < 4000:
                return self._call_ollama(prompt)

            # Pour les grands documents, découper en sections logiques
            sections = self._split_document_into_sections(content)
            
            # Traiter chaque section avec le contexte complet
            results = []
            for i, section in enumerate(sections, 1):
                logger.info(f"Traitement section {i}/{len(sections)}")
                
                section_prompt = f"""Personnalise cette section:

Section:
{section}

Variables:
{vars_json}

Conserve la mise en forme.
"""
                
                # Appeler l'API pour personnaliser la section
                section_result = self._call_ollama(section_prompt)
                
                # Si erreur, utiliser la section originale
                if section_result.startswith("Erreur:") or "Timeout" in section_result:
                    logger.warning(f"Échec de la personnalisation de la section {i}, utilisation de l'original")
                    results.append(section)
                else:
                    results.append(section_result)
            
            # Joindre les sections en préservant les sauts de ligne
            return "\n\n".join(results)
            
        except Exception as e:
            logger.error(f"Erreur lors de la personnalisation: {e}")
            return self._fallback_personalization(content, variables)

    def _call_ollama(self, prompt: str, timeout: int = 5) -> str:  # Timeout réduit à 5s
        try:
            # Limiter la taille du prompt
            if len(prompt) > 1000:  # Réduit de 2000 à 1000
                prompt = prompt[:1000]
            
            # Options simplifiées
            options = {
                "temperature": 0.1,  # Réduit pour plus de rapidité
                "num_predict": 50,   # Réduit de 100 à 50
                "stop": ["\n\n", "###"]
            }
            
            # Appel simplifié à Ollama avec le modèle Mistral
            response = requests.post(
                "http://localhost:11434/api/generate",
                json={"model": "mistral", "prompt": prompt, **options},
                timeout=timeout
            )
            
            if response.status_code == 200:
                return response.json().get("response", "")
            return ""
            
        except Exception as e:
            logger.error(f"Erreur Ollama: {str(e)}")
            return ""

    def _analyze_document_by_sections(self, content: str) -> Dict[str, Any]:
        # Analyse simplifiée
        try:
            sections = self._split_document_into_sections(content)
            results = []
            
            for section in sections[:5]:  # Limite à 5 sections
                if len(section) > 500:  # Réduit de 800 à 500
                    section = section[:500]
                result = self._analyze_section_simplified(section)
                if result:
                    results.append(result)
            
            return self._merge_results(results)
                    
        except Exception as e:
            logger.error(f"Erreur analyse: {str(e)}")
            return {}

    def _analyze_section_simplified(self, content: str) -> Dict[str, Any]:
        try:
            prompt = f"""
Tu es un expert en extraction de variables personnalisables dans les documents. Extrais tous les champs qui devraient être remplacés par des informations du client.

Exemples de variables personnalisables à identifier:
1. Informations personnelles:
   - nom, prénom, civilité (M./Mme/etc.)
   - société, entreprise, organisation
   - fonction, poste, titre professionnel

2. Coordonnées:
   - adresse complète, rue, numéro, ville, code_postal, pays
   - email, téléphone, mobile, fax
   - site_web

3. Données financières:
   - montant, prix, total, sous_total
   - devise, taux_tva, montant_tva
   - reduction, remise, acompte
   - prix_unitaire, quantité

4. Dates et références:
   - date (toute date dans le document)
   - date_creation, date_emission, date_echeance
   - reference, numero_commande, numero_client
   - numero_facture, identifiant

5. Contenu contextuel:
   - objet, sujet, titre_document
   - description, details, motif
   - lieu, emplacement
   - conditions, termes

Ne retourne que les variables qui sont clairement identifiables comme des champs à personnaliser.
Assure-toi que chaque variable soit correctement nommée et catégorisée.

Analyse le texte suivant et identifie toutes les variables personnalisables:
-----
{content}
-----

Réponds uniquement avec un objet JSON de la forme:
{{
  "variables": {{
    "nom_variable": {{
      "valeur": "Valeur actuelle extraite du document",
      "type": "Type de variable (nom, date, montant, etc.)",
      "description": "Description brève de la variable"
    }},
    ...
  }}
}}
"""
            response = self._call_ollama(prompt, timeout=15)
            
            if not response:
                return {}
                
            try:
                # Essayer d'extraire la partie JSON
                json_match = re.search(r'({.*})', response.replace('\n', ''), re.DOTALL)
                if json_match:
                    json_str = json_match.group(1)
                    try:
                        extracted_data = json.loads(json_str)
                        # Vérifier que c'est un dictionnaire et qu'il contient "variables"
                        if isinstance(extracted_data, dict):
                            if "variables" in extracted_data:
                                return extracted_data
                            return {"variables": extracted_data}
                    except:
                        pass
                
                # Si l'extraction a échoué, essayer de nettoyer la réponse entière
                try:
                    return json.loads(response)
                except:
                    pass
                    
                # En dernier recours, parser manuellement
                variable_pattern = r'"([^"]+)":\s*"([^"]+)"'
                matches = re.findall(variable_pattern, response)
                if matches:
                    variables = {}
                    for key, value in matches:
                        variables[key] = value
                    return {"variables": variables}
                    
                return {}
            except Exception as e:
                logger.error(f"Erreur lors du parsing de la réponse JSON: {e}")
                return {}
                        
        except Exception as e:
            logger.error(f"Erreur section: {str(e)}")
            return {}

    def _merge_results(self, results: List[Dict]) -> Dict[str, Any]:
        merged = {}
        for result in results:
            for key, value in result.items():
                if key not in merged:
                    merged[key] = value
        return merged

    def _analyze_fallback(self, content: str) -> Dict[str, Any]:
        """
        Méthode de secours pour l'analyse de document lorsque l'analyse principale échoue
        
        Args:
            content: Contenu textuel du document à analyser
            
        Returns:
            Dict contenant les informations extraites
        """
        logger.info("Utilisation de la méthode de secours pour l'analyse")
        try:
            # Utiliser un prompt simplifié pour Mistral avec moins de complexité
            prompt = f"""
Analyse ce document et identifie tous les champs personnalisables qui devraient être remplacés par les informations d'un client.
Ne retourne qu'un objet JSON simple sans explication.

Document:
---
{content[:2000]}  # Limiter à 2000 caractères pour garantir une réponse
---

Format de réponse attendu:
{{
  "variables": {{
    "nom_variable": "valeur_actuelle",
    ...
  }}
}}
"""
            # Essayer d'abord avec Mistral
            try:
                response = self._call_ollama(prompt, timeout=10)
                if response:
                    # Tenter d'extraire le JSON
                    json_match = re.search(r'({.*})', response.replace('\n', ''), re.DOTALL)
                    if json_match:
                        try:
                            extracted_data = json.loads(json_match.group(1))
                            if isinstance(extracted_data, dict) and "variables" in extracted_data:
                                logger.info(f"Analyse de secours via Mistral réussie, {len(extracted_data['variables'])} variables extraites")
                                return extracted_data
                        except:
                            logger.warning("Erreur d'analyse JSON de la réponse Mistral")
            except Exception as e:
                logger.warning(f"Erreur lors de l'utilisation de Mistral pour l'analyse de secours: {e}")
            
            # Si Mistral échoue, utiliser une analyse par expressions régulières
            variables = {}
            
            # Rechercher des dates (format JJ/MM/AAAA ou JJ-MM-AAAA ou AAAA-MM-JJ)
            dates = re.findall(r'\b\d{1,2}[/.-]\d{1,2}[/.-]\d{2,4}\b|\b\d{4}-\d{2}-\d{2}\b', content)
            if dates:
                variables['date'] = dates[0]
                # Chercher plus de contexte pour identifier le type de date
                date_contexts = {
                    'date_emission': ['émis', 'établi', 'création', 'facture du'],
                    'date_echeance': ['échéance', 'paiement', 'règlement', 'due', 'limite'],
                    'date_livraison': ['livraison', 'livré', 'expédition']
                }
                
                content_lower = content.lower()
                for date in dates[:3]:  # Limiter aux 3 premières dates
                    for date_type, contexts in date_contexts.items():
                        # Chercher le contexte avant et après la date
                        for context in contexts:
                            if re.search(f"{context}.{{0,30}}{re.escape(date)}", content_lower) or \
                               re.search(f"{re.escape(date)}.{{0,30}}{context}", content_lower):
                                variables[date_type] = date
                                break
            
            # Rechercher des montants (avec ou sans devise)
            amount_patterns = [
                r'\b\d+(?:[,.]\d{1,3})?(?:\s?[€$£]|EUR|USD|GBP)\b',  # Montants avec devise
                r'\b\d+(?:[,.]\d{2})(?:\s?HT|\s?TTC)?\b',  # Montants avec HT/TTC
                r'total\s*:?\s*\d+(?:[,.]\d{1,3})?'  # Totaux
            ]
            
            for pattern in amount_patterns:
                amounts = re.findall(pattern, content, re.IGNORECASE)
                if amounts:
                    # Identifier le contexte pour chaque montant
                    amount_contexts = {
                        'montant_total': ['total', 'net à payer', 'montant'],
                        'montant_ht': ['HT', 'hors taxe'],
                        'montant_ttc': ['TTC', 'toutes taxes', 'TVA incluse'],
                        'acompte': ['acompte', 'avance', 'versement initial']
                    }
                    
                    content_lower = content.lower()
                    for amount in amounts[:5]:  # Limiter aux 5 premiers montants
                        amount_clean = amount.strip()
                        for amount_type, contexts in amount_contexts.items():
                            for context in contexts:
                                if re.search(f"{context}.{{0,50}}{re.escape(amount_clean)}", content_lower, re.IGNORECASE) or \
                                   re.search(f"{re.escape(amount_clean)}.{{0,50}}{context}", content_lower, re.IGNORECASE):
                                    variables[amount_type] = amount_clean
                                    break
                                    
                    # Si aucun contexte spécifique n'a été trouvé pour le premier montant
                    if 'montant_total' not in variables and 'montant_ht' not in variables and 'montant_ttc' not in variables:
                        variables['montant'] = amounts[0]
            
            # Rechercher des emails
            emails = re.findall(r'\b[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+\b', content)
            if emails:
                variables['email'] = emails[0]
                
                # Vérifier s'il y a plusieurs emails et les classifier
                if len(emails) > 1:
                    email_contexts = {
                        'email_contact': ['contact', 'info', 'renseignement'],
                        'email_facturation': ['facture', 'facturation', 'compta', 'finance'],
                        'email_support': ['support', 'aide', 'assistance', 'sav']
                    }
                    
                    content_lower = content.lower()
                    for i, email in enumerate(emails[:3]):  # Limiter aux 3 premiers emails
                        for email_type, contexts in email_contexts.items():
                            for context in contexts:
                                if re.search(f"{context}.{{0,50}}{re.escape(email)}", content_lower, re.IGNORECASE) or \
                                   re.search(f"{re.escape(email)}.{{0,50}}{context}", content_lower, re.IGNORECASE):
                                    variables[email_type] = email
                                    break
            
            # Rechercher des numéros de téléphone (formats internationaux inclus)
            phone_patterns = [
                r'\b(?:0\d[\s.-]?){4,5}\d{2}\b',  # Format français
                r'\b\+\d{2}[\s.-]?(?:\d[\s.-]?){8,11}\b',  # Format international
                r'\b(?:\d{2}[\s.-]?){4,5}\b'  # Format général
            ]
            
            for pattern in phone_patterns:
                phones = re.findall(pattern, content)
                if phones:
                    variables['telephone'] = phones[0]
                    # Identifier différents types de téléphones
                    if len(phones) > 1:
                        phone_contexts = {
                            'telephone_mobile': ['mobile', 'portable', 'cell'],
                            'telephone_fixe': ['fixe', 'bureau'],
                            'fax': ['fax', 'télécopie']
                        }
                        
                        content_lower = content.lower()
                        for i, phone in enumerate(phones[:3]):  # Limiter aux 3 premiers numéros
                            for phone_type, contexts in phone_contexts.items():
                                for context in contexts:
                                    if re.search(f"{context}.{{0,30}}{re.escape(phone)}", content_lower, re.IGNORECASE) or \
                                       re.search(f"{re.escape(phone)}.{{0,30}}{context}", content_lower, re.IGNORECASE):
                                        variables[phone_type] = phone
                                        break
                    break
            
            # Rechercher des adresses (codes postaux et villes)
            # Format français: code postal à 5 chiffres suivi de ville
            cp_city_matches = re.findall(r'\b(\d{5})\s+([A-Z][A-Za-zÀ-ÿ\s-]{2,})\b', content)
            if cp_city_matches:
                cp, city = cp_city_matches[0]
                variables['code_postal'] = cp
                variables['ville'] = city
                
                # Chercher une rue associée (avant le code postal)
                content_parts = content.split(cp)
                if len(content_parts) > 1:
                    before_cp = content_parts[0]
                    # Chercher la dernière ligne contenant "rue", "avenue", etc.
                    address_indicators = ['rue', 'avenue', 'boulevard', 'place', 'chemin', 'allée', 'impasse', 'route']
                    for indicator in address_indicators:
                        addr_match = re.findall(f'\\b\\d+[\\s,]*{indicator}[\\sA-Za-zÀ-ÿ,\\-\']+', before_cp, re.IGNORECASE)
                        if addr_match:
                            variables['adresse'] = addr_match[-1].strip()
                            break
            
            # Essayer d'identifier le type de document
            doc_types = {
                'contrat': ['contrat', 'convention', 'accord', 'engagement'],
                'facture': ['facture', 'paiement', 'règlement', 'montant', 'total', 'tva'],
                'devis': ['devis', 'estimation', 'proposition', 'offre'],
                'lettre': ['madame', 'monsieur', 'cordialement', 'sincères salutations'],
                'rapport': ['rapport', 'analyse', 'étude', 'synthèse'],
                'attestation': ['atteste', 'certification', 'confirme', 'attester'],
                'cv': ['expérience', 'compétences', 'formation', 'diplôme', 'curriculum vitae']
            }
            
            content_lower = content.lower()
            for doc_type, keywords in doc_types.items():
                if any(keyword in content_lower for keyword in keywords):
                    variables['type_document'] = doc_type
                    break
            else:
                variables['type_document'] = 'document'
            
            # Rechercher des références de document
            ref_patterns = [
                r'\b(?:ref|référence|n°)(?:\s|:)+([a-zA-Z0-9-_/]+)',
                r'\b(?:facture|devis|commande|bon)(?:\s|:)+n°\s*([a-zA-Z0-9-_/]+)',
                r'\b(?:identifiant|id|numéro)(?:\s|:)+([a-zA-Z0-9-_/]+)'
            ]
            
            for pattern in ref_patterns:
                refs = re.findall(pattern, content, re.IGNORECASE)
                if refs:
                    variables['reference'] = refs[0].strip()
                    break
            
            # Rechercher un nom ou une entité
            name_patterns = [
                r'M(?:\.|onsieur)\s+([A-Z][a-zA-ZéèêëàâäôöûüçÉÈÊËÀÂÄÔÖÛÜÇ\s-]+)',
                r'Mme(?:\.|adame)\s+([A-Z][a-zA-ZéèêëàâäôöûüçÉÈÊËÀÂÄÔÖÛÜÇ\s-]+)',
                r'(?:Société|SARL|SAS|SA|EURL|SASU)\s+([A-Z][a-zA-ZéèêëàâäôöûüçÉÈÊËÀÂÄÔÖÛÜÇ\s-]+)'
            ]
            
            for pattern in name_patterns:
                matches = re.findall(pattern, content)
                if matches:
                    if 'Monsieur' in pattern:
                        variables['civilite'] = 'M.'
                        variables['nom'] = matches[0].strip()
                    elif 'Madame' in pattern:
                        variables['civilite'] = 'Mme'
                        variables['nom'] = matches[0].strip()
                    else:
                        variables['societe'] = matches[0].strip()
                    break
            
            # Si aucun nom ou société, chercher des mots commençant par majuscule
            if 'nom' not in variables and 'societe' not in variables:
                cap_words_pattern = r'(?<=[^\.\?\!]\s)[A-Z][a-zA-ZéèêëàâäôöûüçÉÈÊËÀÂÄÔÖÛÜÇ]{3,}'
                cap_words = re.findall(cap_words_pattern, content)
                if cap_words:
                    # Prendre les mots qui apparaissent au moins deux fois
                    from collections import Counter
                    word_counts = Counter(cap_words)
                    for word, count in word_counts.most_common(3):
                        if count >= 2 and len(word) > 3:  # Au moins 2 occurrences et 4 caractères
                            if re.match(r'^[A-Z][a-z]+$', word):  # Format typique pour un nom propre
                                variables['nom'] = word
                                break
            
            # Créer un résultat structuré
            structured_variables = {}
            for key, value in variables.items():
                structured_variables[key] = {
                    "valeur": value,
                    "type": self._get_variable_type(key),
                    "description": self._get_variable_description(key)
                }
            
            # Créer le résultat final
            result = {
                'variables': structured_variables,
                'detection_method': 'fallback_enhanced',
                'confidence': 'medium'
            }
            
            logger.info(f"Analyse de secours terminée, {len(variables)} variables extraites")
            return result
            
        except Exception as e:
            logger.error(f"Erreur lors de l'analyse de secours: {e}")
            # Retourner un résultat minimal mais valide
            return {
                'variables': {
                    'type_document': {
                        "valeur": "document",
                        "type": "categorie",
                        "description": "Type de document"
                    }
                },
                'detection_method': 'fallback_error',
                'error': str(e)
            }

    def _get_variable_type(self, key: str) -> str:
        """
        Détermine le type d'une variable en fonction de sa clé
        
        Args:
            key: Nom de la variable
            
        Returns:
            Type de la variable
        """
        type_mapping = {
            'nom': 'identite', 'prenom': 'identite', 'civilite': 'identite',
            'societe': 'organisation', 'entreprise': 'organisation', 
            'adresse': 'coordonnees', 'ville': 'coordonnees', 'code_postal': 'coordonnees', 'pays': 'coordonnees',
            'email': 'contact', 'telephone': 'contact', 'telephone_mobile': 'contact', 'telephone_fixe': 'contact',
            'date': 'temporel', 'date_emission': 'temporel', 'date_echeance': 'temporel', 'date_livraison': 'temporel',
            'montant': 'financier', 'montant_ht': 'financier', 'montant_ttc': 'financier', 'tva': 'financier',
            'reference': 'identifiant', 'numero_facture': 'identifiant', 'numero_client': 'identifiant',
            'type_document': 'categorie', 'objet': 'contenu', 'description': 'contenu'
        }
        return type_mapping.get(key, 'autre')
    
    def _get_variable_description(self, key: str) -> str:
        """
        Fournit une description pour une variable en fonction de sa clé
        
        Args:
            key: Nom de la variable
            
        Returns:
            Description de la variable
        """
        description_mapping = {
            'nom': 'Nom de famille ou nom complet',
            'prenom': 'Prénom de la personne',
            'civilite': 'Civilité (M., Mme, etc.)',
            'societe': 'Nom de la société ou organisation',
            'entreprise': 'Nom de l\'entreprise',
            'adresse': 'Adresse postale',
            'ville': 'Ville de l\'adresse',
            'code_postal': 'Code postal',
            'pays': 'Pays de l\'adresse',
            'email': 'Adresse email principale',
            'telephone': 'Numéro de téléphone',
            'telephone_mobile': 'Numéro de téléphone mobile',
            'telephone_fixe': 'Numéro de téléphone fixe',
            'date': 'Date générique mentionnée dans le document',
            'date_emission': 'Date d\'émission du document',
            'date_echeance': 'Date d\'échéance de paiement',
            'date_livraison': 'Date de livraison prévue',
            'montant': 'Montant global',
            'montant_ht': 'Montant hors taxes',
            'montant_ttc': 'Montant toutes taxes comprises',
            'tva': 'Montant ou taux de TVA',
            'reference': 'Référence du document',
            'numero_facture': 'Numéro de facture',
            'numero_client': 'Numéro ou identifiant client',
            'type_document': 'Type ou catégorie du document',
            'objet': 'Objet ou sujet du document',
            'description': 'Description ou détails'
        }
        return description_mapping.get(key, f'Valeur de {key}')

    def diagnose_analysis_issues(self, file_path: str) -> Dict[str, Any]:
        """
        Diagnostique les problèmes potentiels dans l'analyse d'un document
        
        Args:
            file_path: Chemin vers le document à analyser
            
        Returns:
            Dict contenant les diagnostics et informations sur les problèmes
        """
        logger.info(f"Diagnostic de l'analyse pour le fichier: {file_path}")
        
        diagnostics = {
            "file_path": file_path,
            "file_exists": os.path.exists(file_path),
            "file_size": os.path.getsize(file_path) if os.path.exists(file_path) else 0,
            "file_extension": os.path.splitext(file_path)[1].lower() if os.path.exists(file_path) else "",
            "text_extraction": None,
            "text_length": 0,
            "mistral_connection": False,
            "mistral_response": None,
            "issues": [],
            "recommendations": []
        }
        
        # Vérification du fichier
        if not diagnostics["file_exists"]:
            diagnostics["issues"].append("Le fichier n'existe pas")
            diagnostics["recommendations"].append("Vérifiez le chemin du fichier")
            return diagnostics
            
        if diagnostics["file_size"] == 0:
            diagnostics["issues"].append("Le fichier est vide")
            diagnostics["recommendations"].append("Vérifiez que le fichier contient des données")
            return diagnostics
            
        # Extraction de texte
        try:
            content = self.convert_to_text(file_path)
            diagnostics["text_extraction"] = "success"
            diagnostics["text_length"] = len(content)
            
            # Échantillon du texte
            diagnostics["text_sample"] = content[:200] + "..." if len(content) > 200 else content
            
            if diagnostics["text_length"] == 0:
                diagnostics["issues"].append("Aucun texte n'a pu être extrait du document")
                diagnostics["recommendations"].append("Vérifiez que le document contient du texte et pas seulement des images")
                
        except Exception as e:
            diagnostics["text_extraction"] = "failed"
            diagnostics["issues"].append(f"Erreur lors de l'extraction du texte: {str(e)}")
            diagnostics["recommendations"].append("Vérifiez le format du document et réessayez avec un autre format si possible")
            
        # Test de la connexion à Mistral
        try:
            test_prompt = "Renvoie simplement le mot 'OK'"
            response = self._call_ollama(test_prompt, timeout=3)
            
            diagnostics["mistral_connection"] = "OK" in response
            diagnostics["mistral_response"] = response[:50] if response else "Pas de réponse"
            
            if not diagnostics["mistral_connection"]:
                diagnostics["issues"].append("Problème de connexion avec le modèle Mistral")
                diagnostics["recommendations"].append("Vérifiez que le serveur Ollama est en cours d'exécution et que le modèle Mistral est installé")
                
        except Exception as e:
            diagnostics["issues"].append(f"Erreur lors du test de connexion à Mistral: {str(e)}")
            diagnostics["recommendations"].append("Vérifiez la configuration du serveur Ollama et assurez-vous qu'il fonctionne correctement")
            
        # Analyse conditionnelle du contenu
        if diagnostics["text_extraction"] == "success" and diagnostics["text_length"] > 0:
            # Vérification pour les formats structurés/semi-structurés
            if diagnostics["file_extension"] in [".pdf", ".docx", ".odt"]:
                if content.count("\n") < 5:
                    diagnostics["issues"].append("Le document semble manquer de structure (peu de sauts de ligne)")
                    diagnostics["recommendations"].append("Vérifiez que la mise en forme du document est préservée")
            
            # Vérification des caractères spéciaux/problèmes d'encodage
            strange_chars_count = len(re.findall(r'[^\x00-\x7F\u00C0-\u00FF]', content))
            if strange_chars_count > len(content) * 0.1:  # Plus de 10% de caractères étranges
                diagnostics["issues"].append("Le document contient beaucoup de caractères non standard")
                diagnostics["recommendations"].append("Vérifiez l'encodage du document")
            
        # Ajout d'un statut global
        if not diagnostics["issues"]:
            diagnostics["status"] = "ok"
            diagnostics["recommendations"].append("Aucun problème détecté, l'analyse devrait fonctionner correctement")
        else:
            diagnostics["status"] = "issues_detected"
            
        logger.info(f"Diagnostic terminé, statut: {diagnostics['status']}")
        return diagnostics

    def _split_document_into_sections(self, content: str) -> List[str]:
        """
        Divise un document en sections logiques pour analyse
        
        Args:
            content: Contenu du document à diviser
            
        Returns:
            Liste des sections du document
        """
        try:
            logger.info("Division du document en sections pour analyse")
            
            # Si le contenu est très court, le traiter comme une seule section
            if len(content) < 1000:
                return [content]
                
            # Tentative de division basée sur les titres et les sauts de ligne
            sections = []
            
            # Identifier les marqueurs de séparation de sections
            section_markers = [
                # Titres numérotés (1. Introduction, 1.1 Contexte, etc.)
                r'\n\s*(?:\d+\.)+\s*[A-Z][^\n]{3,}',
                # Titres en majuscules
                r'\n\s*[A-Z][A-Z\s]{5,}[A-Z]\s*\n',
                # Titres avec formatage spécial (encadrés par des astérisques, soulignés, etc.)
                r'\n\s*[\*\-=_]{2,}[^\n]{3,}[\*\-=_]{2,}',
                # Séparateurs horizontaux
                r'\n\s*[\-=_\*]{3,}\s*\n',
                # Doubles sauts de ligne suivis d'un mot en majuscule
                r'\n\s*\n\s*[A-Z][^\n]{3,}',
                # Mots-clés spécifiques qui indiquent souvent un changement de section
                r'\n(?:ARTICLE|ANNEXE|SECTION|CHAPITRE|PARTIE)[^\n]*'
            ]
            
            # Combiner tous les marqueurs en une seule expression régulière
            combined_pattern = '|'.join(section_markers)
            
            # Trouver tous les points de séparation
            split_points = []
            for m in re.finditer(combined_pattern, content):
                split_points.append(m.start())
                
            # Ajouter le début et la fin du document
            split_points = [0] + split_points + [len(content)]
            
            # Créer les sections
            for i in range(len(split_points) - 1):
                start = split_points[i]
                end = split_points[i + 1]
                section_content = content[start:end].strip()
                
                # Ne pas ajouter de sections vides
                if section_content:
                    sections.append(section_content)
                    
            # Si aucune section n'a été identifiée, diviser le document en paragraphes
            if len(sections) <= 1:
                paragraphs = re.split(r'\n\s*\n', content)
                
                # Regrouper les paragraphes en sections d'environ 1000 caractères
                current_section = ""
                for para in paragraphs:
                    if para.strip():
                        if len(current_section) + len(para) > 1000:
                            if current_section:
                                sections.append(current_section)
                                current_section = para
                            else:
                                sections.append(para)
                        else:
                            current_section += "\n\n" + para if current_section else para
                
                # Ajouter la dernière section
                if current_section:
                    sections.append(current_section)
            
            # Si le document est très long mais qu'on n'a pas trouvé beaucoup de sections,
            # diviser en morceaux de taille similaire
            if len(content) > 5000 and len(sections) < 3:
                # Réinitialiser les sections
                sections = []
                
                # Diviser en sections de 1000 caractères environ
                section_size = 1000
                for i in range(0, len(content), section_size):
                    section = content[i:i + section_size]
                    if section.strip():
                        sections.append(section)
                
            logger.info(f"Document divisé en {len(sections)} sections")
            return sections
                
        except Exception as e:
            logger.error(f"Erreur lors de la division du document en sections: {e}")
            # En cas d'erreur, retourner le document entier comme une seule section
            return [content]

    def _fallback_personalization(self, content: str, variables: Dict[str, Any]) -> str:
        """
        Méthode de personnalisation de secours pour remplacer les variables par leurs valeurs
        
        Args:
            content: Contenu du document à personnaliser
            variables: Dictionnaire des variables à remplacer
            
        Returns:
            Document personnalisé
        """
        logger.info("Utilisation de la méthode de personnalisation de secours")
        try:
            result = content
            
            # Préparer les variables pour le remplacement
            replace_dict = {}
            for var_name, var_info in variables.items():
                value = ""
                
                # Extraire la valeur selon le format (dict ou valeur simple)
                if isinstance(var_info, dict):
                    if "valeur" in var_info:
                        value = var_info["valeur"]
                    elif "current_value" in var_info:
                        value = var_info["current_value"]
                else:
                    value = str(var_info)
                    
                if value:
                    # Créer différentes variantes de recherche pour le nom de variable
                    patterns = [
                        r'\{' + re.escape(var_name) + r'\}',  # {nom_variable}
                        r'\{\{' + re.escape(var_name) + r'\}\}',  # {{nom_variable}}
                        r'\[\[' + re.escape(var_name) + r'\]\]',  # [[nom_variable]]
                        r'<' + re.escape(var_name) + r'>',  # <nom_variable>
                        r'%' + re.escape(var_name) + r'%',  # %nom_variable%
                        r'\$' + re.escape(var_name) + r'\$',  # $nom_variable$
                        r'\$\{' + re.escape(var_name) + r'\}',  # ${nom_variable}
                        r'\b' + re.escape(var_name) + r'\b'  # nom_variable (mot entier)
                    ]
                    
                    # Ajouter aussi la version avec des underscores remplacés par des espaces
                    if '_' in var_name:
                        space_name = var_name.replace('_', ' ')
                        patterns.append(r'\b' + re.escape(space_name) + r'\b')
                    
                    # Pour certaines variables spécifiques, ajouter des variantes courantes
                    special_cases = {
                        'nom': ['NOM', 'Nom'],
                        'prenom': ['PRENOM', 'Prénom', 'PRÉNOM'],
                        'adresse': ['ADRESSE', 'Adresse'],
                        'ville': ['VILLE', 'Ville'],
                        'code_postal': ['CODE POSTAL', 'Code postal', 'CP', 'cp'],
                        'date': ['DATE', 'Date']
                    }
                    
                    if var_name.lower() in special_cases:
                        for variant in special_cases[var_name.lower()]:
                            patterns.append(r'\b' + re.escape(variant) + r'\b')
                    
                    # Appliquer tous les motifs de recherche
                    for pattern in patterns:
                        result = re.sub(pattern, value, result)
            
            # Recherche avancée pour les variables non explicitement marquées
            # Par exemple, pour capturer "Nom: ________" ou "Adresse: ______"
            common_fields = {
                'nom': [r'Nom\s*:\s*_{3,}', r'Nom\s*:\s*\.{3,}', r'Nom\s*:(?:\s*$|\s{10,})'],
                'prenom': [r'Pr[ée]nom\s*:\s*_{3,}', r'Pr[ée]nom\s*:\s*\.{3,}', r'Pr[ée]nom\s*:(?:\s*$|\s{10,})'],
                'adresse': [r'Adresse\s*:\s*_{3,}', r'Adresse\s*:\s*\.{3,}', r'Adresse\s*:(?:\s*$|\s{10,})'],
                'ville': [r'Ville\s*:\s*_{3,}', r'Ville\s*:\s*\.{3,}', r'Ville\s*:(?:\s*$|\s{10,})'],
                'code_postal': [r'Code postal\s*:\s*_{3,}', r'Code postal\s*:\s*\.{3,}', r'Code postal\s*:(?:\s*$|\s{10,})'],
                'email': [r'E-?mail\s*:\s*_{3,}', r'E-?mail\s*:\s*\.{3,}', r'E-?mail\s*:(?:\s*$|\s{10,})'],
                'telephone': [r'T[ée]l[ée]phone\s*:\s*_{3,}', r'T[ée]l[ée]phone\s*:\s*\.{3,}', r'T[ée]l[ée]phone\s*:(?:\s*$|\s{10,})']
            }
            
            for field, patterns in common_fields.items():
                # Vérifier si nous avons une valeur pour ce champ
                value = ""
                for var_name, var_info in variables.items():
                    if var_name.lower() == field or var_name.lower().endswith('_' + field):
                        if isinstance(var_info, dict):
                            if "valeur" in var_info:
                                value = var_info["valeur"]
                            elif "current_value" in var_info:
                                value = var_info["current_value"]
                        else:
                            value = str(var_info)
                        break
                
                if value:
                    for pattern in patterns:
                        result = re.sub(pattern, f"{field.capitalize()}: {value}", result, flags=re.IGNORECASE)
            
            logger.info("Personnalisation de secours terminée")
            return result
            
        except Exception as e:
            logger.error(f"Erreur lors de la personnalisation de secours: {e}")
            # En cas d'erreur, retourner le document original
            return content
